"""Paste the 13 audit screenshots for a SKU into a target PPTX.

Slide layout assumptions (matches ClaudeMattel.pptx as of May 2026):
  - Slide 14 → 2 ChatGPT snaps (chatgpt_initial, chatgpt)
  - Slide 15 → 2 Google snaps (initial, full_page)
  - Slide 24 → 3 Alexa pill snaps (alexa_top, alexa_inline, alexa_qa)
  - Slide 25 → 3 Alexa response snaps (likes, dislikes, certs)
  - Slide 26 → 3 Alexa response snaps (materials, alternatives, budget_alt)

Images are placed in a single row along the bottom edge of each target
slide so they don't cover existing content. Existing slide content is
not modified; only new pictures are added. The user crops/repositions
manually in PowerPoint.

Usage:
    python src/paste_to_pptx.py --sku <name> --pptx <path>
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent.parent

# Map 0-based slide index -> ordered list of screenshot filenames.
# Change here if the template's slide structure shifts.
#
# Layout per ClaudeMattel.pptx structure (May 2026):
#   13 = Google         14 = ChatGPT
#   21 = Sparky chat    22 = Walmart 1-star reviews
#   23 = Alexa pills    24, 25 = Alexa Q+A
LAYOUT: dict[int, list[str]] = {
    13: ["chatgpt_initial.png", "chatgpt.png"],
    14: ["initial.png", "full_page.png"],
    21: ["sparky_likes.png", "sparky_dislikes.png", "sparky_defects.png"],
    22: ["walmart_bad_review_1.png", "walmart_bad_review_2.png",
         "walmart_bad_review_3.png"],
    23: ["alexa_top.png", "alexa_inline.png", "alexa_qa.png"],
    24: ["alexa_likes.png", "alexa_dislikes.png", "alexa_certs.png"],
    25: ["alexa_materials.png", "alexa_alternatives.png", "alexa_budget_alt.png"],
}

# Per-engine file groups for the `--engine X --slide N` mode used by
# the per-engine slash commands (/gemini, /gpt, /alexa, /sparky,
# /walmart). Each entry is a LIST of file groups; each group becomes
# one consecutive slide starting at the user-specified slide.
ENGINE_GROUPS: dict[str, list[list[str]]] = {
    "google": [["initial.png", "full_page.png"]],
    "gemini": [["initial.png", "full_page.png"]],
    "chatgpt": [["chatgpt_initial.png", "chatgpt.png"]],
    "gpt": [["chatgpt_initial.png", "chatgpt.png"]],
    "alexa": [
        ["alexa_top.png", "alexa_inline.png", "alexa_qa.png"],
        ["alexa_likes.png", "alexa_dislikes.png", "alexa_certs.png"],
        ["alexa_materials.png", "alexa_alternatives.png", "alexa_budget_alt.png"],
    ],
    "sparky": [["sparky_likes.png", "sparky_dislikes.png", "sparky_defects.png"]],
    "walmart": [["walmart_bad_review_1.png", "walmart_bad_review_2.png",
                 "walmart_bad_review_3.png"]],
}


def paste(sku: str, pptx_path: Path,
          only_slides: set[int] | None = None) -> int:
    samples = ROOT / "samples" / sku
    if not samples.is_dir():
        print(f"ERROR: samples dir not found: {samples}", file=sys.stderr)
        return 1

    if not pptx_path.is_file():
        print(f"ERROR: pptx not found: {pptx_path}", file=sys.stderr)
        return 1

    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    print(f"PPTX:    {pptx_path}")
    print(f"Samples: {samples}")
    print(f"Slide size: {slide_w/914400:.2f}x{slide_h/914400:.2f} in\n")

    margin = Inches(0.1)
    gap = Inches(0.1)
    max_h = Inches(3.0)

    placed_total = 0
    for slide_idx, filenames in LAYOUT.items():
        if only_slides is not None and slide_idx not in only_slides:
            continue
        if slide_idx >= len(prs.slides):
            print(f"WARN: slide {slide_idx + 1} doesn't exist in deck — skipping",
                  file=sys.stderr)
            continue
        slide = prs.slides[slide_idx]

        n = len(filenames)
        available_w = slide_w - 2 * margin - (n - 1) * gap
        per_w = available_w // n
        print(f"Slide {slide_idx + 1}: placing {n} image(s)")
        for i, fname in enumerate(filenames):
            path = samples / fname
            if not path.is_file():
                print(f"  ! missing {fname} — skipping", file=sys.stderr)
                continue
            im = PILImage.open(path)
            iw, ih = im.size
            aspect = ih / iw
            target_w = per_w
            target_h = int(target_w * aspect)
            if target_h > max_h:
                target_h = max_h
                target_w = int(target_h / aspect)
            x = margin + i * (per_w + gap) + (per_w - target_w) // 2
            y = slide_h - target_h - Inches(0.1)
            slide.shapes.add_picture(str(path), x, y, target_w, target_h)
            print(f"  + {fname}  {target_w/914400:.2f}x{target_h/914400:.2f}in")
            placed_total += 1

    prs.save(str(pptx_path))
    print(f"\nSaved -> {pptx_path}")
    print(f"Pasted {placed_total} image(s) total.")
    return 0


def paste_engine(sku: str, pptx_path: Path, engine: str,
                 start_slide: int) -> int:
    """Paste one engine's screenshots onto consecutive slides starting
    at `start_slide` (1-based). Used by the per-engine slash commands.
    Auto-creates the PPTX if it doesn't exist (with enough blank
    slides to land the paste)."""
    engine = engine.lower()
    if engine not in ENGINE_GROUPS:
        print(f"ERROR: unknown engine {engine!r}. Choices: "
              f"{sorted(ENGINE_GROUPS)}", file=sys.stderr)
        return 1
    groups = ENGINE_GROUPS[engine]
    samples = ROOT / "samples" / sku
    if not samples.is_dir():
        print(f"ERROR: samples dir not found: {samples}", file=sys.stderr)
        return 1

    # If the PPTX doesn't exist, create a blank deck with enough slides.
    if not pptx_path.is_file():
        pptx_path.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        # Add slides up to start_slide + len(groups) - 1
        needed = start_slide + len(groups) - 1
        blank_layout = prs.slide_layouts[6]  # blank layout
        while len(prs.slides) < needed:
            prs.slides.add_slide(blank_layout)
        prs.save(str(pptx_path))
        print(f"Created new PPTX with {needed} blank slides at {pptx_path}")

    prs = Presentation(str(pptx_path))
    # Add slides if the deck doesn't have enough.
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 \
        else prs.slide_layouts[-1]
    needed = start_slide + len(groups) - 1
    while len(prs.slides) < needed:
        prs.slides.add_slide(blank_layout)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    print(f"PPTX:    {pptx_path}")
    print(f"Samples: {samples}")
    print(f"Engine:  {engine} ({len(groups)} slide(s) starting at "
          f"slide {start_slide})")
    print(f"Slide size: {slide_w/914400:.2f}x{slide_h/914400:.2f} in\n")

    margin = Inches(0.1)
    gap = Inches(0.1)
    max_h = Inches(5.0)

    placed_total = 0
    for offset, filenames in enumerate(groups):
        slide_num = start_slide + offset  # 1-based
        slide_idx = slide_num - 1
        slide = prs.slides[slide_idx]
        n = len(filenames)
        available_w = slide_w - 2 * margin - (n - 1) * gap
        per_w = available_w // n
        print(f"Slide {slide_num}: placing {n} image(s)")
        for i, fname in enumerate(filenames):
            path = samples / fname
            if not path.is_file():
                print(f"  ! missing {fname} — skipping", file=sys.stderr)
                continue
            im = PILImage.open(path)
            iw, ih = im.size
            aspect = ih / iw
            target_w = per_w
            target_h = int(target_w * aspect)
            if target_h > max_h:
                target_h = max_h
                target_w = int(target_h / aspect)
            x = margin + i * (per_w + gap) + (per_w - target_w) // 2
            y = slide_h - target_h - Inches(0.1)
            slide.shapes.add_picture(str(path), x, y, target_w, target_h)
            print(f"  + {fname}  {target_w/914400:.2f}x{target_h/914400:.2f}in")
            placed_total += 1

    prs.save(str(pptx_path))
    print(f"\nSaved -> {pptx_path}")
    print(f"Pasted {placed_total} image(s) total.")
    return 0


def main() -> int:
    p = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    p.add_argument("--sku", required=True,
                   help="Sample folder name under samples/")
    p.add_argument("--pptx", required=True, type=Path,
                   help="Target .pptx file (overwritten in place)")
    p.add_argument("--only-slides", default=None,
                   help="Comma-separated 1-based slide numbers to paste "
                        "(e.g. '14,15'). Default: all slides in LAYOUT.")
    p.add_argument("--engine", default=None,
                   choices=sorted(ENGINE_GROUPS),
                   help="Single-engine mode: paste just this engine's "
                        "files starting at --slide. Auto-creates the "
                        "PPTX if it doesn't exist.")
    p.add_argument("--slide", type=int, default=None,
                   help="1-based starting slide number for --engine mode.")
    args = p.parse_args()
    if args.engine:
        if args.slide is None:
            print("ERROR: --engine requires --slide", file=sys.stderr)
            return 2
        return paste_engine(args.sku, args.pptx, args.engine, args.slide)
    only = None
    if args.only_slides:
        # User passes 1-based slide numbers; LAYOUT uses 0-based indices.
        only = {int(s.strip()) - 1 for s in args.only_slides.split(",")
                if s.strip()}
    return paste(args.sku, args.pptx, only_slides=only)


if __name__ == "__main__":
    sys.exit(main())
